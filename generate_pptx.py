"""Bold dark-gradient PPT with vibrant colors and detailed content."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Vibrant palette
ELECTRIC = RGBColor(0x7C,0x3A,0xED)  # vivid purple
HOT     = RGBColor(0xEC,0x48,0x99)   # hot pink
CYAN    = RGBColor(0x22,0xD3,0xEE)   # bright cyan
LIME    = RGBColor(0x84,0xCC,0x16)   # lime
ORANGE  = RGBColor(0xF9,0x73,0x16)   # orange
DEEP    = RGBColor(0xF8,0xFA,0xFC)   # light bg
DARK2   = RGBColor(0xFF,0xFF,0xFF)   # card bg
CARD    = RGBColor(0xFF,0xFF,0xFF)   # card
WHITE   = RGBColor(0xFF,0xFF,0xFF)
SLATE   = RGBColor(0x1E,0x29,0x3B)   # dark text
LGRAY   = RGBColor(0x7C,0x3A,0xED)   # accent text
MGRAY   = RGBColor(0x47,0x55,0x69)   # body text
BDR     = RGBColor(0xE2,0xE8,0xF0)   # border
SW = Inches(13.333); SH = Inches(7.5)

def _bg(sl, c=DEEP):
    f = sl.background.fill; f.solid(); f.fore_color.rgb = c

def _rect(sl, l, t, w, h, c):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background(); return s

def _rrect(sl, l, t, w, h, c, bdr=None):
    s = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = c
    if bdr: s.line.color.rgb = bdr; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def _circ(sl, l, t, sz, c):
    s = sl.shapes.add_shape(MSO_SHAPE.OVAL, l, t, sz, sz)
    s.fill.solid(); s.fill.fore_color.rgb = c; s.line.fill.background(); return s

def _tb(sl, l, t, w, h, txt, sz=18, b=False, c=WHITE, al=PP_ALIGN.LEFT):
    tb = sl.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz); p.font.bold = b
    p.font.color.rgb = c; p.font.name = "Segoe UI"; p.alignment = al; return tf

def _bul(sl, l, t, w, h, items, sz=15, c=MGRAY):
    tb = sl.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = it; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.name = "Segoe UI"; p.space_after = Pt(8)

def _hdr(sl, title, accent=ELECTRIC):
    _rect(sl, 0, 0, Inches(0.12), SH, accent)
    _rect(sl, 0, 0, SW, Inches(0.04), RGBColor(0xE2,0xE8,0xF0))
    _tb(sl, Inches(0.5), Inches(0.3), Inches(10), Inches(0.8), title, sz=34, b=True, c=SLATE)
    _rect(sl, Inches(0.5), Inches(1.1), Inches(2), Inches(0.05), accent)

# ═══ SLIDES ═══

def s1(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, RGBColor(0x1E,0x29,0x3B))
    _circ(sl, Inches(9), Inches(-1), Inches(5), RGBColor(0x7C,0x3A,0xED))
    _circ(sl, Inches(10.5), Inches(3), Inches(4), RGBColor(0xEC,0x48,0x99))
    _circ(sl, Inches(-2), Inches(5), Inches(4), RGBColor(0x22,0xD3,0xEE))
    _rect(sl, Inches(1.2), Inches(2.2), Inches(3), Inches(0.07), CYAN)
    _tb(sl, Inches(1.2), Inches(2.5), Inches(9), Inches(1.2),
        "[Your Title Here]", sz=48, b=True, c=WHITE)
    _tb(sl, Inches(1.2), Inches(4.0), Inches(8), Inches(0.6),
        "[Subtitle / Team Name]", sz=22, c=RGBColor(0xA7,0x8B,0xFA))
    _tb(sl, Inches(1.2), Inches(4.8), Inches(8), Inches(0.5),
        "[Department · University · Date]", sz=16, c=MGRAY)
    _tb(sl, Inches(1.2), Inches(5.6), Inches(8), Inches(0.5),
        "[Add additional info]", sz=14, c=RGBColor(0x6B,0x72,0x80))

def s2(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _hdr(sl, "Introduction", CYAN)
    _tb(sl, Inches(0.5), Inches(1.5), Inches(12), Inches(1),
        "The Automated Assignment Evaluation System is a full-stack web platform "
        "developed for NIST University that digitizes and automates the entire "
        "assignment lifecycle — from creation to grading to feedback delivery.",
        sz=17, c=MGRAY)
    cards = [
        (CYAN, "🎯 Objective", "Eliminate manual grading bottlenecks\nby automating evaluation of both\ncoding and descriptive assignments\nusing NLP and sandboxed execution"),
        (HOT,  "👥 Target Users", "Faculty members who create and\nmanage assignments, and students\nwho submit work and receive\ndetailed performance feedback"),
        (ELECTRIC, "⚡ Key Value", "Instant, objective, and consistent\ngrading with built-in peer-to-peer\nplagiarism detection — no external\ntools or licenses required"),
    ]
    for i, (clr, title, desc) in enumerate(cards):
        l = Inches(0.3 + i * 4.2)
        _rrect(sl, l, Inches(3.0), Inches(3.9), Inches(4), CARD, bdr=BDR)
        _rect(sl, l, Inches(3.0), Inches(3.9), Inches(0.06), clr)
        _tb(sl, l+Inches(0.3), Inches(3.3), Inches(3.3), Inches(0.5), title, sz=17, b=True, c=clr)
        _tb(sl, l+Inches(0.3), Inches(4.0), Inches(3.3), Inches(2.5), desc, sz=14, c=MGRAY)

def s3(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _hdr(sl, "Problem Statement", HOT)
    probs = [
        (HOT, "⏱", "Time-Consuming Evaluation",
         "Faculty spend 3-5 days grading a single\nassignment batch of 100+ students.\nThis delays feedback and slows the\nlearning loop significantly."),
        (ORANGE, "🎯", "Subjective & Inconsistent",
         "Different evaluators apply varying\nstandards to the same submission.\nStudents receive unfair grades based\non who evaluates their work."),
        (CYAN, "💸", "Expensive Plagiarism Tools",
         "University-licensed plagiarism checkers\ncost thousands annually. Smaller\ndepartments often go without any\nplagiarism detection at all."),
        (ELECTRIC, "📉", "Zero Data Insights",
         "No centralized analytics exist to track\nclass-wide performance trends,\nidentify struggling students, or\nmeasure assignment difficulty."),
    ]
    for i, (clr, icon, title, desc) in enumerate(probs):
        l = Inches(0.3 + i * 3.2)
        _rrect(sl, l, Inches(1.6), Inches(3.0), Inches(5.5), CARD, bdr=BDR)
        _circ(sl, l+Inches(0.85), Inches(1.9), Inches(1.2), clr)
        _tb(sl, l+Inches(0.85), Inches(2.0), Inches(1.2), Inches(1), icon, sz=36, al=PP_ALIGN.CENTER, c=WHITE)
        _tb(sl, l+Inches(0.15), Inches(3.4), Inches(2.7), Inches(0.5), title, sz=14, b=True, c=clr, al=PP_ALIGN.CENTER)
        _tb(sl, l+Inches(0.15), Inches(4.0), Inches(2.7), Inches(3), desc, sz=12, c=MGRAY, al=PP_ALIGN.CENTER)

def s4(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _hdr(sl, "The Solution", LIME)
    _tb(sl, Inches(0.5), Inches(1.5), Inches(12), Inches(1),
        "A single, unified platform that handles the complete assignment lifecycle — "
        "from faculty-defined rubrics and test cases, through student file uploads, to "
        "automated NLP and code evaluation, peer plagiarism scanning, and detailed grade reports.",
        sz=17, c=MGRAY)
    pts = [
        (LIME, "Auto-grade coding assignments in Python, Java, and C++ by compiling, executing "
               "against faculty-defined test cases, and comparing stdout output"),
        (CYAN, "Evaluate descriptive essays using cosine similarity against reference answers, "
               "plus spelling, grammar, and clarity analysis via TextBlob and NLTK"),
        (HOT,  "Detect plagiarism across all peer submissions using a dual-method approach: "
               "N-gram word overlap (60% weight) and TF-IDF cosine similarity (40% weight)"),
        (ORANGE,"Provide faculty with real-time analytics dashboards showing score distributions, "
                "class averages, plagiarism risk heatmaps, and performance trend charts"),
        (ELECTRIC,"Give students instant access to detailed feedback reports including rubric "
                  "breakdowns, strengths, weaknesses, and improvement suggestions"),
    ]
    for i, (clr, text) in enumerate(pts):
        t = Inches(2.8 + i * 0.9)
        _rect(sl, Inches(0.5), t, Inches(0.12), Inches(0.65), clr)
        _rrect(sl, Inches(0.8), t, Inches(11.7), Inches(0.65), CARD, bdr=BDR)
        _tb(sl, Inches(1.1), t+Inches(0.08), Inches(11.2), Inches(0.5), text, sz=13, c=MGRAY)

def s5(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _hdr(sl, "Core Features", ELECTRIC)
    feats = [
        (CYAN, "🤖","Sandboxed Code Runner",
         "Compiles and executes Python, Java,\nC++ in isolated subprocesses with\nstrict 5-second timeouts to prevent\ninfinite loops or resource abuse"),
        (HOT, "📝","NLP Essay Grader",
         "Compares student essays against\nfaculty reference using TF-IDF cosine\nsimilarity. Checks spelling, grammar,\nclarity, and text completeness"),
        (ELECTRIC, "🔍","Plagiarism Detection",
         "Peer-to-peer comparison engine using\ndual methods: N-gram word overlap\n(60%) + TF-IDF cosine similarity (40%)\nClamped to 0-100% accuracy"),
        (ORANGE, "📊","Analytics Dashboard",
         "Real-time charts showing score\ndistribution, class averages, plagiarism\nrisk levels, and performance trends\nvia Chart.js visualizations"),
        (LIME, "📤","Batch Evaluation",
         "One-click processing of all pending\nsubmissions. Faculty can review\nauto-generated feedback before\nreleasing grades to students"),
        (RGBColor(0xA7,0x8B,0xFA), "🔐","Role-Based Access",
         "Separate faculty and student portals.\nStudents can only view their own\nreports. Faculty control grade release\ntiming and visibility"),
    ]
    for i, (clr, icon, title, desc) in enumerate(feats):
        row, col = i % 2, i // 2
        l = Inches(0.3 + col * 4.2); t = Inches(1.5 + row * 2.9)
        _rrect(sl, l, t, Inches(4.0), Inches(2.6), CARD, bdr=BDR)
        _rect(sl, l, t, Inches(4.0), Inches(0.06), clr)
        _circ(sl, l+Inches(0.2), t+Inches(0.25), Inches(0.65), clr)
        _tb(sl, l+Inches(0.2), t+Inches(0.3), Inches(0.65), Inches(0.5), icon, sz=20, al=PP_ALIGN.CENTER, c=WHITE)
        _tb(sl, l+Inches(1.05), t+Inches(0.3), Inches(2.7), Inches(0.4), title, sz=15, b=True, c=clr)
        _tb(sl, l+Inches(0.2), t+Inches(0.9), Inches(3.6), Inches(1.5), desc, sz=12, c=MGRAY)

def s6(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _hdr(sl, "Tech Stack", ORANGE)
    items = [
        (ELECTRIC, "Python 3.x", "Core backend language powering all services"),
        (HOT,      "Flask + Blueprints", "Lightweight MVC framework with modular routing"),
        (CYAN,     "SQLite + SQLAlchemy", "Embedded database with full ORM for models"),
        (LIME,     "TextBlob / NLTK", "NLP libraries for sentiment, spelling, similarity"),
        (ORANGE,   "Flask-Login", "Session-based auth with role-based access control"),
        (RGBColor(0xA7,0x8B,0xFA), "Chart.js", "Interactive charts for dashboard analytics"),
        (CYAN,     "HTML5 / CSS3 / JS", "Responsive UI with custom CSS design system"),
        (HOT,      "Subprocess Module", "Sandboxed code compilation and execution"),
    ]
    for i, (clr, name, desc) in enumerate(items):
        row, col = i % 4, i // 4
        l = Inches(0.3 + col * 6.4); t = Inches(1.5 + row * 1.4)
        _rrect(sl, l, t, Inches(6.1), Inches(1.15), CARD, bdr=BDR)
        _rect(sl, l, t, Inches(0.1), Inches(1.15), clr)
        _tb(sl, l+Inches(0.3), t+Inches(0.1), Inches(3.5), Inches(0.4), name, sz=16, b=True, c=clr)
        _tb(sl, l+Inches(0.3), t+Inches(0.5), Inches(5.5), Inches(0.5), desc, sz=13, c=MGRAY)

def s7(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _hdr(sl, "System Workflow", CYAN)
    steps = [
        (ELECTRIC,"1","CREATE","Faculty defines assignment\ntype, rubric weights,\nreference solution & test cases"),
        (HOT,    "2","SUBMIT","Students upload code files\nor essays in PDF/TXT format\nvia the submission portal"),
        (LIME,   "3","PARSE", "File parser extracts raw text\nfrom uploaded documents\nwith error-safe handling"),
        (ORANGE, "4","GRADE", "Code Runner executes against\ntest cases or NLP Engine\nevaluates against reference"),
        (CYAN,   "5","DETECT","Plagiarism engine compares\nevery submission against\nall peers using dual methods"),
        (RGBColor(0xA7,0x8B,0xFA),"6","REPORT","Final scores aggregated from\nrubric weights. Detailed\nfeedback report generated"),
    ]
    _rect(sl, Inches(0.8), Inches(3.7), Inches(11.8), Inches(0.04), BDR)
    for i, (clr, num, title, desc) in enumerate(steps):
        l = Inches(0.2 + i * 2.15)
        dot = _circ(sl, l+Inches(0.55), Inches(3.45), Inches(0.55), clr)
        tf = dot.text_frame; tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(18); tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = WHITE; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _tb(sl, l, Inches(2.3), Inches(1.9), Inches(0.8), title, sz=13, c=clr, al=PP_ALIGN.CENTER, b=True)
        _tb(sl, l, Inches(4.2), Inches(1.9), Inches(1.8), desc, sz=11, c=MGRAY, al=PP_ALIGN.CENTER)
        if i < 5:
            _tb(sl, l+Inches(1.7), Inches(3.55), Inches(0.5), Inches(0.4), "→", sz=20, c=RGBColor(0x3B,0x2A,0x5C), al=PP_ALIGN.CENTER)
    _rrect(sl, Inches(2), Inches(6.2), Inches(9.3), Inches(0.8), CARD, bdr=RGBColor(0x3B,0x2A,0x5C))
    _tb(sl, Inches(2.3), Inches(6.3), Inches(8.7), Inches(0.6),
        "Batch Evaluation: One click processes ALL pending submissions simultaneously", sz=14, c=LGRAY, al=PP_ALIGN.CENTER)

def s8(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _hdr(sl, "Database Schema", ELECTRIC)
    _tb(sl, Inches(0.5), Inches(1.4), Inches(12), Inches(0.4),
        "4 core models connected via SQLAlchemy ORM on SQLite:", sz=15, c=MGRAY)
    models = [
        (CYAN, "User", "id  ·  username  ·  email\npassword_hash  ·  role\nfull_name  ·  student_id\n\nRoles: faculty | student\nSecure bcrypt hashing"),
        (ELECTRIC, "Assignment", "id  ·  title  ·  description\ntype  ·  language  ·  due_date\nreference_solution\ntest_cases (JSON)\n\nTypes: coding | descriptive"),
        (HOT, "Submission", "id  ·  student_id (FK)\nassignment_id (FK)\nfile_path  ·  raw_text\nstatus  ·  submitted_at\n\nStatus: pending | evaluated"),
        (ORANGE, "Evaluation", "id  ·  submission_id (FK)\nscore  ·  letter_grade\nfeedback_json  ·  sub_scores\nplagiarism_score\nreleased (boolean)"),
    ]
    for i, (clr, name, fields) in enumerate(models):
        l = Inches(0.2 + i * 3.25)
        _rrect(sl, l, Inches(2.0), Inches(3.1), Inches(4.8), CARD, bdr=RGBColor(0x3B,0x2A,0x5C))
        _rect(sl, l, Inches(2.0), Inches(3.1), Inches(0.55), clr)
        _tb(sl, l+Inches(0.3), Inches(2.1), Inches(2.5), Inches(0.4), name, sz=18, b=True, c=WHITE)
        _tb(sl, l+Inches(0.3), Inches(2.7), Inches(2.5), Inches(3.8), fields, sz=12, c=MGRAY)
    _tb(sl, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.3),
        "User ──1:N──▶ Submission ──1:1──▶ Evaluation     │     Assignment ──1:N──▶ Submission",
        sz=12, b=True, c=LGRAY, al=PP_ALIGN.CENTER)

def s9(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _hdr(sl, "Security & Integrity", HOT)
    items = [
        (CYAN, "🔐","Role-Based Access Control",
         "Strict route decorators ensure students cannot access faculty\ndashboards, evaluation triggers, or peer submission reports"),
        (HOT, "📁","16 MB Upload Limit",
         "MAX_CONTENT_LENGTH enforced at the Flask layer with a custom\n413 error handler to prevent disk exhaustion attacks"),
        (ELECTRIC, "⏱","5-Second Execution Timeout",
         "All student code runs in subprocess with strict timeout.\nInfinite loops are killed instantly, returning clean error messages"),
        (ORANGE, "✅","Rubric Weight Validation",
         "System validates rubric weights sum to 100%. If not, it auto-\nnormalizes them during evaluation to prevent score inflation"),
        (LIME, "📝","Production Logging",
         "RotatingFileHandler writes to instance/app.log with automatic\nrotation at 10MB. Tracks all errors, logins, and evaluations"),
        (RGBColor(0xA7,0x8B,0xFA), "🛡","Sanitized Error Output",
         "Compiler errors and runtime exceptions are caught and returned\nas helpful feedback — no raw stack traces exposed to users"),
    ]
    for i, (clr, icon, title, desc) in enumerate(items):
        row, col = i % 3, i // 3
        l = Inches(0.3 + col * 6.4); t = Inches(1.5 + row * 1.9)
        _rrect(sl, l, t, Inches(6.1), Inches(1.65), CARD, bdr=RGBColor(0x3B,0x2A,0x5C))
        _circ(sl, l+Inches(0.2), t+Inches(0.3), Inches(0.7), clr)
        _tb(sl, l+Inches(0.2), t+Inches(0.35), Inches(0.7), Inches(0.5), icon, sz=22, al=PP_ALIGN.CENTER, c=WHITE)
        _tb(sl, l+Inches(1.1), t+Inches(0.15), Inches(4.8), Inches(0.4), title, sz=14, b=True, c=clr)
        _tb(sl, l+Inches(1.1), t+Inches(0.55), Inches(4.8), Inches(0.9), desc, sz=12, c=MGRAY)

def s10(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl); _hdr(sl, "Impact & Future Scope", LIME)
    _tb(sl, Inches(0.5), Inches(1.4), Inches(5.8), Inches(0.5), "Current Impact", sz=20, b=True, c=LIME)
    imps = [
        "✔  Reduces grading time from 3-5 days to under 5 minutes",
        "✔  Ensures 100% consistent, objective evaluation across sections",
        "✔  Eliminates need for expensive third-party plagiarism tools",
        "✔  Students receive instant, detailed feedback upon grade release",
        "✔  Faculty gain data-driven class performance insights",
    ]
    _bul(sl, Inches(0.5), Inches(2.0), Inches(5.8), Inches(4), imps, sz=14, c=MGRAY)
    _tb(sl, Inches(7), Inches(1.4), Inches(5.8), Inches(0.5), "Future Roadmap", sz=20, b=True, c=ELECTRIC)
    futs = [
        (CYAN,     "🐳  Docker Containers", "Ephemeral sandboxes for secure code execution"),
        (HOT,      "🧠  LLM Integration",   "Ollama-hosted models for semantic essay feedback"),
        (ORANGE,   "⚡  Async Queues",       "Celery + Redis for non-blocking batch grading"),
        (ELECTRIC, "🔗  LMS Integration",    "REST API hooks to sync with university systems"),
    ]
    for i, (clr, title, desc) in enumerate(futs):
        t = Inches(2.1 + i * 1.2)
        _rect(sl, Inches(7), t, Inches(0.1), Inches(0.9), clr)
        _rrect(sl, Inches(7.2), t, Inches(5.6), Inches(0.9), CARD, bdr=RGBColor(0x3B,0x2A,0x5C))
        _tb(sl, Inches(7.5), t+Inches(0.08), Inches(5), Inches(0.35), title, sz=14, b=True, c=clr)
        _tb(sl, Inches(7.5), t+Inches(0.42), Inches(5), Inches(0.35), desc, sz=12, c=MGRAY)

def s11(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6]); _bg(sl, RGBColor(0x1E,0x29,0x3B))
    _circ(sl, Inches(9), Inches(-1), Inches(5), ELECTRIC)
    _circ(sl, Inches(10.5), Inches(3), Inches(4), HOT)
    _circ(sl, Inches(-2), Inches(5), Inches(4), CYAN)
    _rect(sl, Inches(4.5), Inches(2.0), Inches(4.3), Inches(0.07), CYAN)
    _tb(sl, Inches(1), Inches(2.3), Inches(11.3), Inches(1),
        "Conclusion", sz=44, b=True, c=WHITE, al=PP_ALIGN.CENTER)
    _tb(sl, Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.5),
        "The Automated Assignment Evaluation System transforms how NIST University "
        "handles academic assessment. By combining sandboxed code execution, NLP-driven "
        "essay grading, and peer plagiarism detection into one integrated platform, it "
        "delivers faster, fairer, and more insightful grading — empowering both faculty "
        "and students to focus on what matters most: learning.",
        sz=16, c=RGBColor(0x9C,0xA3,0xAF), al=PP_ALIGN.CENTER)
    _tb(sl, Inches(1), Inches(5.5), Inches(11.3), Inches(0.5),
        "Thank You  ·  Questions & Discussion", sz=22, b=True, c=CYAN, al=PP_ALIGN.CENTER)


def main():
    prs = Presentation(); prs.slide_width = SW; prs.slide_height = SH
    for fn in [s1,s2,s3,s4,s5,s6,s7,s8,s9,s10,s11]: fn(prs)
    out = "NIST_AutoEval_Presentation_v3.pptx"
    prs.save(out); print(f"Saved: {out} ({len(prs.slides)} slides)")

if __name__ == "__main__": main()
